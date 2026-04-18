"""
Genera el pitch de SunHub en PPTX (16:9).
Paleta corporativa Stitch/SunHub:
  primary        #16A34A
  primary-dark   #006B2C
  accent         #FACC15
  info           #0EA5E9
  surface        #F4FCF0
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pathlib import Path

# ---- Paleta ----
PRIMARY      = RGBColor(0x16, 0xA3, 0x4A)
PRIMARY_DARK = RGBColor(0x00, 0x6B, 0x2C)
ACCENT       = RGBColor(0xFA, 0xCC, 0x15)
INFO         = RGBColor(0x0E, 0xA5, 0xE9)
SURFACE      = RGBColor(0xF4, 0xFC, 0xF0)
INK          = RGBColor(0x0F, 0x17, 0x2A)
INK_SOFT     = RGBColor(0x33, 0x41, 0x55)
MUTED        = RGBColor(0x64, 0x74, 0x8B)
BORDER       = RGBColor(0xE2, 0xE8, 0xF0)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
DANGER       = RGBColor(0xDC, 0x26, 0x26)
ROSE_BG      = RGBColor(0xFE, 0xF2, 0xF2)
ROSE_BORDER  = RGBColor(0xFE, 0xCA, 0xCA)
EMERALD_BG   = SURFACE
EMERALD_BORD = RGBColor(0xBB, 0xF7, 0xD0)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]


# ---- helpers --------------------------------------------------------------
def add_rect(slide, x, y, w, h, fill, line=None, line_w=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        if line_w: shp.line.width = line_w
    shp.shadow.inherit = False
    return shp

def add_rounded(slide, x, y, w, h, fill, line=None, line_w=None, radius=0.05):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.adjustments[0] = radius
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        if line_w: shp.line.width = line_w
    shp.shadow.inherit = False
    return shp

def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=INK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.text = ""
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb

def add_bullets(slide, x, y, w, h, items, *, size=16, color=INK_SOFT,
                line_spacing=1.25, bullet_color=None):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    bc = bullet_color or PRIMARY
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(6)
        bullet = p.add_run()
        bullet.text = "• "
        bullet.font.name = "Calibri"
        bullet.font.size = Pt(size)
        bullet.font.bold = True
        bullet.font.color.rgb = bc
        run = p.add_run()
        run.text = item
        run.font.name = "Calibri"
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return tb

def add_circle(slide, cx, cy, r, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx - r, cy - r, 2*r, 2*r)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp

def slide_frame(slide, *, page_label=None, page_total=6, kicker=None):
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), PRIMARY)
    add_text(slide, Inches(0.6), Inches(7.08), Inches(6), Inches(0.3),
             "SunHub · Techos Rentables · Hackathon Tinku 2026",
             size=9, color=MUTED)
    if page_label is not None:
        add_text(slide, Inches(11.9), Inches(7.08), Inches(1.1), Inches(0.3),
                 f"{page_label} / {page_total}", size=9, color=MUTED, align=PP_ALIGN.RIGHT)
    if kicker:
        add_rounded(slide, Inches(0.6), Inches(0.55), Inches(2.6), Inches(0.36),
                    SURFACE, line=PRIMARY, line_w=Pt(0.75), radius=0.5)
        add_text(slide, Inches(0.6), Inches(0.55), Inches(2.6), Inches(0.36),
                 kicker, size=10, bold=True, color=PRIMARY_DARK,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ==========================================================================
# Slide 1 — Portada
# ==========================================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, PRIMARY_DARK)
add_rect(s, Inches(0), Inches(0), SLIDE_W, Inches(4.6), PRIMARY)
add_circle(s, Inches(11.6), Inches(1.6), Inches(0.9), ACCENT)
add_circle(s, Inches(12.2), Inches(2.3), Inches(0.35), WHITE)
add_text(s, Inches(0.9), Inches(0.8), Inches(6), Inches(0.5),
         "SUNHUB", size=14, bold=True, color=ACCENT)
add_text(s, Inches(0.9), Inches(2.2), Inches(11.5), Inches(1.8),
         "Una sola plataforma\npara toda tu operación solar.",
         size=54, bold=True, color=WHITE)
add_text(s, Inches(0.9), Inches(4.9), Inches(11.5), Inches(1),
         "Centraliza datos, anticipa fallas y cumple contratos —\n"
         "de Huawei, Growatt, Hoymiles, Deye y más, en tiempo real.",
         size=20, color=WHITE)
add_rounded(s, Inches(0.9), Inches(6.3), Inches(5.2), Inches(0.55), ACCENT, radius=0.5)
add_text(s, Inches(0.9), Inches(6.3), Inches(5.2), Inches(0.55),
         "Techos Rentables · Equipo Los Incapaces",
         size=12, bold=True, color=INK, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(0.9), Inches(7.08), Inches(12), Inches(0.3),
         "1 / 6", size=9, color=WHITE, align=PP_ALIGN.RIGHT)

# ==========================================================================
# Slide 2 — Problema
# ==========================================================================
s = prs.slides.add_slide(BLANK)
slide_frame(s, page_label=2, kicker="EL PROBLEMA")
add_text(s, Inches(0.6), Inches(1.1), Inches(12), Inches(0.9),
         "Datos fragmentados. Operación manual. Riesgo contractual.",
         size=30, bold=True, color=INK)
add_text(s, Inches(0.6), Inches(1.95), Inches(12), Inches(0.6),
         "200+ plantas · 6+ fabricantes · cero visibilidad unificada.",
         size=16, color=MUTED)

cards = [
    ("130+ h/mes", "en trabajo manual recurrente", INFO),
    ("~40 min", "para generar un reporte por cliente", ACCENT),
    (">5 min", "para detectar una falla — llega tarde al SLA", DANGER),
    ("Riesgo", "de penalización por incumplimiento", PRIMARY_DARK),
]
card_w = Inches(2.85); card_h = Inches(2.1)
gap = Inches(0.2); base_x = Inches(0.6); base_y = Inches(3.0)
for i, (kpi, label, col) in enumerate(cards):
    x = base_x + (card_w + gap) * i
    add_rounded(s, x, base_y, card_w, card_h, WHITE, line=BORDER, line_w=Pt(0.75), radius=0.08)
    add_rect(s, x, base_y, Inches(0.08), card_h, col)
    add_text(s, x + Inches(0.35), base_y + Inches(0.35), card_w - Inches(0.5), Inches(0.9),
             kpi, size=30, bold=True, color=col)
    add_text(s, x + Inches(0.35), base_y + Inches(1.2), card_w - Inches(0.5), Inches(0.9),
             label, size=13, color=INK_SOFT)

add_text(s, Inches(0.6), Inches(5.6), Inches(12), Inches(0.9),
         "Consecuencia: mantenimiento reactivo, operación no escalable y margen erosionado.",
         size=14, color=INK_SOFT)

# ==========================================================================
# Slide 3 — Solución
# ==========================================================================
s = prs.slides.add_slide(BLANK)
slide_frame(s, page_label=3, kicker="LA SOLUCIÓN")
add_text(s, Inches(0.6), Inches(1.1), Inches(12), Inches(0.9),
         "Un ecosistema que unifica, predice y actúa.",
         size=30, bold=True, color=INK)
add_text(s, Inches(0.6), Inches(1.95), Inches(12), Inches(0.5),
         "Cinco capas. Una sola experiencia.",
         size=15, color=MUTED)

steps = [
    ("1", "API central",           "Huawei, Growatt, Deye, Hoymiles\nunificados en un modelo canónico", INFO),
    ("2", "Predicción",            "IA + reglas de negocio\nanticipan fallas y degradación",           ACCENT),
    ("3", "Mantenimiento ágil",    "Aceptar, ticketizar y remediar\nsin salir del sistema",            PRIMARY),
    ("4", "Clima inteligente",     "Pronóstico + score para elegir\nel día ideal de mantenimiento",   INFO),
    ("5", "App unificada",         "Operador, cliente y reportes\nautomáticos en una sola vista",     PRIMARY_DARK),
]
sw = Inches(2.3); sh = Inches(3.3)
sg = Inches(0.18); sx = Inches(0.6); sy = Inches(2.80)
for i, (num, t, d, col) in enumerate(steps):
    x = sx + (sw + sg) * i
    add_rounded(s, x, sy, sw, sh, WHITE, line=BORDER, line_w=Pt(0.75), radius=0.06)
    add_circle(s, x + Inches(0.55), sy + Inches(0.55), Inches(0.35), col)
    add_text(s, x, sy + Inches(0.3), Inches(1.1), Inches(0.5),
             num, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(0.2), sy + Inches(1.1), sw - Inches(0.4), Inches(0.8),
             t, size=15, bold=True, color=INK)
    add_text(s, x + Inches(0.2), sy + Inches(1.75), sw - Inches(0.4), Inches(1.5),
             d, size=11, color=INK_SOFT)
    if i < len(steps) - 1:
        arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                 x + sw + Emu(20000), sy + sh/2 - Inches(0.1),
                                 sg - Emu(40000), Inches(0.2))
        arr.fill.solid(); arr.fill.fore_color.rgb = PRIMARY
        arr.line.fill.background()

add_text(s, Inches(0.6), Inches(6.35), Inches(12), Inches(0.4),
         "Loop cerrado: predecir → actuar → medir → aprender.",
         size=13, bold=True, color=PRIMARY_DARK)

# ==========================================================================
# Slide 4 — Transformación (rediseñada: filas pareadas)
# ==========================================================================
s = prs.slides.add_slide(BLANK)
slide_frame(s, page_label=4, kicker="TRANSFORMACIÓN")
add_text(s, Inches(0.6), Inches(1.1), Inches(12), Inches(0.8),
         "Del caos operativo a la operación predictiva.",
         size=30, bold=True, color=INK)
add_text(s, Inches(0.6), Inches(1.92), Inches(12), Inches(0.4),
         "Cada dimensión de la operación, transformada.",
         size=14, color=MUTED)

# Encabezados de tres columnas
head_y = Inches(2.55)
col_x   = Inches(0.6)                 # dimensión (1.8)
antes_x = Inches(2.5)                 # antes    (4.6)
desp_x  = Inches(8.0)                 # después  (4.6)

# cápsulas de cabecera
add_rounded(s, col_x,   head_y, Inches(1.8), Inches(0.42), WHITE)
add_text   (s, col_x,   head_y, Inches(1.8), Inches(0.42),
            "DIMENSIÓN", size=10, bold=True, color=MUTED,
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)

add_rounded(s, antes_x, head_y, Inches(4.9), Inches(0.42),
            ROSE_BG, line=ROSE_BORDER, line_w=Pt(0.5), radius=0.5)
add_text   (s, antes_x, head_y, Inches(4.9), Inches(0.42),
            "ANTES", size=11, bold=True, color=DANGER,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_rounded(s, desp_x,  head_y, Inches(4.7), Inches(0.42),
            EMERALD_BG, line=EMERALD_BORD, line_w=Pt(0.5), radius=0.5)
add_text   (s, desp_x,  head_y, Inches(4.7), Inches(0.42),
            "CON SUNHUB", size=11, bold=True, color=PRIMARY_DARK,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# Filas: (dimensión, antes, después)
rows = [
    ("Datos",          "6+ portales separados por marca",       "1 API canónica · multi-proveedor"),
    ("Reportes",       "~40 min por cliente · manual",          "Segundos · automáticos y agendados"),
    ("Detección",      "Reactiva, después del incidente",       "Alarmas en <5 min · IA prioriza"),
    ("Mantenimiento",  "Visita presencial · decisión a ciegas", "Remoto vía API · día óptimo con clima"),
    ("Decisión",       "Intuición + Excel",                     "Copilot AI con datos en vivo"),
]

row_h = Inches(0.7); row_gap = Inches(0.1); start_y = Inches(3.15)
for i, (dim, before, after) in enumerate(rows):
    y = start_y + (row_h + row_gap) * i
    # columna dimensión (barra amarilla delgada al inicio)
    add_rect(s, col_x, y + Inches(0.12), Inches(0.08), Inches(0.46), ACCENT)
    add_text(s, col_x + Inches(0.2), y, Inches(1.7), row_h,
             dim, size=14, bold=True, color=INK, anchor=MSO_ANCHOR.MIDDLE)
    # antes
    add_rounded(s, antes_x, y, Inches(4.9), row_h, ROSE_BG,
                line=ROSE_BORDER, line_w=Pt(0.5), radius=0.15)
    add_text(s, antes_x + Inches(0.35), y, Inches(4.4), row_h,
             before, size=13, color=INK_SOFT, anchor=MSO_ANCHOR.MIDDLE)
    # flecha entre columnas
    arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                             antes_x + Inches(4.95), y + row_h/2 - Inches(0.09),
                             Inches(0.45), Inches(0.18))
    arr.fill.solid(); arr.fill.fore_color.rgb = PRIMARY
    arr.line.fill.background()
    # después
    add_rounded(s, desp_x, y, Inches(4.7), row_h, EMERALD_BG,
                line=EMERALD_BORD, line_w=Pt(0.5), radius=0.15)
    add_text(s, desp_x + Inches(0.35), y, Inches(4.3), row_h,
             after, size=13, bold=True, color=PRIMARY_DARK, anchor=MSO_ANCHOR.MIDDLE)

# Remate
add_text(s, Inches(0.6), Inches(7.0), Inches(12), Inches(0.3),
         "Misma flota · mismos contratos · otra operación.",
         size=12, bold=True, color=MUTED)

# ==========================================================================
# Slide 5 — Impacto
# ==========================================================================
s = prs.slides.add_slide(BLANK)
slide_frame(s, page_label=5, kicker="IMPACTO")
add_text(s, Inches(0.6), Inches(1.1), Inches(12), Inches(0.9),
         "Más tiempo, menos riesgo, más cumplimiento.",
         size=30, bold=True, color=INK)
add_text(s, Inches(0.6), Inches(1.95), Inches(12), Inches(0.5),
         "Los números hablan solos.",
         size=15, color=MUTED)

big = [
    ("De 40 min\na segundos", "Reporte por cliente", PRIMARY),
    ("<5 min", "Detección de fallas",            INFO),
    ("130+ h/mes", "Liberadas del trabajo manual", ACCENT),
]
bw = Inches(4.0); bh = Inches(2.3); by = Inches(2.7); gx = Inches(0.15)
for i, (kpi, label, col) in enumerate(big):
    x = Inches(0.6) + (bw + gx) * i
    add_rounded(s, x, by, bw, bh, WHITE, line=BORDER, line_w=Pt(0.75), radius=0.06)
    add_rect(s, x, by, bw, Inches(0.15), col)
    add_text(s, x, by + Inches(0.5), bw, Inches(1.2),
             kpi, size=26, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_text(s, x, by + Inches(1.7), bw, Inches(0.5),
             label, size=13, color=INK_SOFT, align=PP_ALIGN.CENTER)

soft = [
    ("Centralización total",     "1 plataforma · 6 marcas · 200+ plantas"),
    ("Operación escalable",      "Sin sumar cabezas cuando crece la flota"),
    ("Cumplimiento contractual", "SLA, PR, uptime y CO₂ bajo control"),
]
for i, (t, d) in enumerate(soft):
    x = Inches(0.6) + (bw + gx) * i
    y = Inches(5.3)
    add_rect(s, x, y, Inches(0.12), Inches(1.3), PRIMARY)
    add_text(s, x + Inches(0.25), y, bw - Inches(0.25), Inches(0.5),
             t, size=14, bold=True, color=INK)
    add_text(s, x + Inches(0.25), y + Inches(0.5), bw - Inches(0.25), Inches(0.9),
             d, size=12, color=INK_SOFT)

# ==========================================================================
# Slide 6 — Cierre
# ==========================================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, PRIMARY_DARK)
add_rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), ACCENT)
add_circle(s, Inches(12.0), Inches(6.2), Inches(1.4), PRIMARY)
add_circle(s, Inches(11.3), Inches(5.7), Inches(0.3), ACCENT)

add_rounded(s, Inches(0.9), Inches(1.1), Inches(2.2), Inches(0.4), ACCENT, radius=0.5)
add_text(s, Inches(0.9), Inches(1.1), Inches(2.2), Inches(0.4),
         "CIERRE", size=11, bold=True, color=INK,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_text(s, Inches(0.9), Inches(2.2), Inches(11.5), Inches(3),
         "No unificamos solo datos.\nUnificamos decisiones, operación\ny cumplimiento.",
         size=52, bold=True, color=WHITE)

add_text(s, Inches(0.9), Inches(5.6), Inches(11.5), Inches(0.6),
         "SunHub — el sistema operativo de tu flota solar.",
         size=20, color=ACCENT)

add_text(s, Inches(0.9), Inches(6.4), Inches(11.5), Inches(0.5),
         "Techos Rentables · Hackathon Tinku 2026 · Equipo Los Incapaces",
         size=12, color=WHITE)
add_text(s, Inches(11.5), Inches(7.08), Inches(1.5), Inches(0.3),
         "6 / 6", size=9, color=WHITE, align=PP_ALIGN.RIGHT)

# ---- Save ----
out = Path(__file__).parent / "SunHub-Pitch.pptx"
prs.save(out)
print(f"OK -> {out}")
